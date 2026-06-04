#!/usr/bin/env python3
"""slides-emit.py — render clinical workflow diagrams + Slides Pattern Library to .pptx.

Two emitters share one set of primitives:
- emit_diagram(spec, path): a workflow diagram from a spec dict
- emit_pl(path): the Slides Pattern Library (one .pptx covering diagram + deck families)

The PL is the canonical surface per the architectural pivot (workflow plan
Decision 2026-06-04h, extended 2026-06-04i): one Slides PL covers BOTH diagram
primitives and deck primitives. Both families share a base layer (palette,
typography, page chrome). Workflow diagrams + future decks are applications of
PL primitives, not bespoke compositions.

Brand-fidelity intent (haven canon — Lab/cena-health-brand/principles/):
- Lora commands (title); Source Sans 3 works (body); Source Code Pro for
  system/identifier labels (lane names, watchers, TBD callouts).
- Sand family: structural register. Warm-amber / red / teal: severity / accent.
- Restraint: brand moments earn attention through spacing, not saturation.

v2/v3 fixes baked into primitives:
- Triangle arrow heads on connectors (OXML <a:tailEnd type="triangle"/>)
- Dominant-axis edge routing (v0.3): horizontal-dominant edges exit right/left
  and enter left/right; vertical-dominant exit bottom/top and enter top/bottom.
  STRAIGHT when the secondary-axis offset is below LINEAR_TOLERANCE_SVG;
  ELBOW (bentConnector3 Z-shape) otherwise.
- Edge label mask sized to text width + small padding (arrow heads stay visible)
- Callout anchored BOTTOM-up so single-line and multi-line callouts hold a
  consistent gap to the node above
- Muted edges render dashed (OXML <a:prstDash val="sysDash"/>) + thicker line
- Connector shadow inheritance disabled (off-brand drop shadow)

Honest limits (v2 spike):
- Connectors are placed by coordinates, not shape-attached. Box drag → connector
  doesn't follow automatically (collaborative iteration still works).
- ELBOW label position is the (x1,y1)-(x2,y2) midpoint; the actual bend may not
  pass through it exactly. Mask still works (sits on bg); placement approximate.
- Slide bounds are real. For complex content that overflows, layout validation
  is Phase B work (plan entry B10 — slides-spec layout validator).
- ELBOW corners render sharp (90°). Rounded corners require custom-geometry
  shapes via <a:custGeom> with lnTo + arcTo path commands — Phase B item B12.
- Font rendering depends on availability in Slides; we set names declaratively
  (Slides substitutes if absent — Lora → Cambria-ish; SS3 → Arial-ish).

Usage:
  slides-emit.py <out.pptx>           # workflow diagram (escalation spec)
  slides-emit.py --pl <out.pptx>      # Slides Pattern Library
"""

import sys
import os
from pathlib import Path

# Ensure spec_from_markdown can be imported when running from any cwd
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ============================================================================
# Shared base layer — palette, typography, geometry
# ============================================================================

# Brand palette (Cena Color System v2)
COLORS = {
    "page_bg":   RGBColor(0xFB, 0xFA, 0xF8),
    "sand_50":   RGBColor(0xF8, 0xF4, 0xEC),
    "sand_100":  RGBColor(0xE4, 0xDF, 0xD7),
    "sand_200":  RGBColor(0xCF, 0xCA, 0xC2),
    "sand_300":  RGBColor(0xBB, 0xB6, 0xAD),
    "sand_500":  RGBColor(0x87, 0x7F, 0x76),
    "sand_600":  RGBColor(0x77, 0x70, 0x69),
    "sand_700":  RGBColor(0x5A, 0x54, 0x4E),
    "sand_900":  RGBColor(0x25, 0x21, 0x1D),
    "amber_50":  RGBColor(0xFA, 0xEF, 0xDA),
    "amber_500": RGBColor(0xAA, 0x82, 0x32),
    "amber_700": RGBColor(0x73, 0x53, 0x11),
    "red_50":    RGBColor(0xFF, 0xE9, 0xE5),
    "red_500":   RGBColor(0xC1, 0x3C, 0x3B),
    "red_700":   RGBColor(0x8A, 0x2A, 0x29),
    "teal_700":  RGBColor(0x1E, 0x51, 0x49),
    "arrow_default": RGBColor(0x8C, 0x7C, 0x5E),
    "arrow_muted":   RGBColor(0xB8, 0xA9, 0x88),
}

MODIFIER_STYLES = {
    "default":              {"fill": "sand_50",  "border": "sand_300", "border_width": 1.0},
    "uncertain-tbd":        {"fill": "amber_50", "border": "amber_500", "border_width": 1.0},
    "uncertain-assumption": {"fill": "sand_50",  "border": "sand_500",  "border_width": 1.0},
    "uncertain-gap":        {"fill": "red_50",   "border": "red_500",   "border_width": 1.5},
    "attestation-gate":     {"fill": "sand_50",  "border": "teal_700",  "border_width": 2.5},
}

CALLOUT_COLORS = {
    "tbd":        "amber_700",
    "critical":   "red_700",
    "gap":        "red_500",
    "assumption": "sand_700",
}

# Geometry: SVG viewBox px → Slide inches
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
LEFT_MARGIN_IN = 0.45
TOP_MARGIN_IN = 0.95
SCALE = 0.0085
LANE_LABEL_X_IN = 0.15
CALLOUT_GAP_TO_NODE_IN = 0.10  # consistent visual gap callout-bottom → node-top
CALLOUT_BOX_H_IN = 0.55        # accommodates 2-line callouts via BOTTOM anchor
LINEAR_TOLERANCE_SVG = 10      # secondary-axis offset below this → STRAIGHT
CORNER_RADIUS_PT = 7           # rounded-corner radius on custom-geometry Z connectors
OVERLAP_OFFSET_SVG = 22        # bend-point perpendicular offset between edges sharing a node pair

# Per-corner arc parameters for custom-geometry rounded Z paths.
# Key: (incoming_direction, outgoing_direction)
# Value: (stAng_deg, swAng_deg) for <a:arcTo>; OOXML uses 60000ths of a degree.
# Angles: east=0°, south=90°, west=180°, north=270° (clockwise from east).
ARC_PARAMS = {
    ("east", "south"):   (270,  90),
    ("east", "north"):   ( 90, -90),
    ("west", "south"):   (270, -90),
    ("west", "north"):   ( 90,  90),
    ("south", "east"):   (180, -90),
    ("south", "west"):   (  0,  90),
    ("north", "east"):   (180,  90),
    ("north", "west"):   (  0, -90),
}

# Shape ID counter for custom-geometry shapes (avoid collision with python-pptx auto-IDs)
_CUSTOM_SHAPE_ID = [10000]


def x_in(x_svg):
    return Inches(LEFT_MARGIN_IN + x_svg * SCALE)


def y_in(y_svg):
    return Inches(TOP_MARGIN_IN + y_svg * SCALE)


def w_in(w_svg):
    return Inches(w_svg * SCALE)


def h_in(h_svg):
    return Inches(h_svg * SCALE)


# --- OXML helpers ---

def _add_arrow_tail(line_format, kind="triangle", width="med", length="med"):
    """Inject <a:tailEnd type='triangle' w='med' len='med'/> into a connector's <a:ln>.

    python-pptx exposes <a:ln> via line_format._get_or_add_ln() but has no helper
    for headEnd/tailEnd; we drop down to lxml.
    """
    ln = line_format._get_or_add_ln()
    for existing in ln.findall(qn("a:tailEnd")):
        ln.remove(existing)
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", kind)
    tail.set("w", width)
    tail.set("len", length)


def _set_line_dash(line_format, val="sysDash"):
    """Inject <a:prstDash val='sysDash'/> into a connector's <a:ln>.

    Schema position: prstDash comes before tailEnd within <a:ln>. If tailEnd
    is already present, insert prstDash before it; otherwise append.
    """
    ln = line_format._get_or_add_ln()
    for existing in ln.findall(qn("a:prstDash")):
        ln.remove(existing)
    dash = etree.Element(qn("a:prstDash"))
    dash.set("val", val)
    tail = ln.find(qn("a:tailEnd"))
    if tail is not None:
        tail.addprevious(dash)
    else:
        ln.append(dash)


def _bbox_overlaps(a, b):
    """Both bboxes are (x_svg, y_svg, w_svg, h_svg). True if rectangles intersect."""
    ax1, ay1, aw, ah = a
    bx1, by1, bw, bh = b
    return not (ax1 + aw <= bx1 or bx1 + bw <= ax1 or
                ay1 + ah <= by1 or by1 + bh <= ay1)


def _collides_any(bbox, others):
    return any(_bbox_overlaps(bbox, o) for o in others)


def _node_bbox(node):
    return (node["x"], node["y"], node["w"], node["h"])


def _rgb_hex(color):
    """RGBColor → 'RRGGBB' uppercase hex (OOXML srgbClr val format)."""
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _add_rounded_z_connector(slide, x1_svg, y1_svg, x2_svg, y2_svg,
                              line_color_rgb, line_width_pt,
                              dash_val=None, r_pt=CORNER_RADIUS_PT,
                              horizontal_first=True,
                              mid_offset_svg=0):
    """Add a Z-shape connector with rounded corners via custom-geometry shape.

    Replaces MSO_CONNECTOR_TYPE.ELBOW (bentConnector3) because the preset's
    bend orientation is determined by its bounding box, not by which shape
    edges the endpoints touch — so it can't be forced into horizontal-first
    routing when we want the last segment to enter perpendicular to a side
    edge. Custom geometry gives us full control + rounded corners.

    horizontal_first=True (default):
        Path: (x1,y1) → (mx,y1) → (mx,y2) → (x2,y2). First and last segments
        horizontal, middle vertical. Arrow at end moves horizontal — perpendicular
        to destination's vertical (left/right) edge.

    horizontal_first=False:
        Path: (x1,y1) → (x1,my) → (x2,my) → (x2,y2). First and last segments
        vertical, middle horizontal. Arrow at end moves vertical — perpendicular
        to destination's horizontal (top/bottom) edge.

    All coordinates in svg-px (consistent with the rest of the module). Internally
    converts to EMU for the custom-geometry path.
    """
    x1_emu = int(x_in(x1_svg).emu)
    y1_emu = int(y_in(y1_svg).emu)
    x2_emu = int(x_in(x2_svg).emu)
    y2_emu = int(y_in(y2_svg).emu)
    r_emu = int(Pt(r_pt).emu)
    line_w_emu = int(Pt(line_width_pt).emu)

    # mid_offset_svg shifts the bend point perpendicular to the dominant axis,
    # so two edges sharing the same undirected node pair don't trace identical
    # Z paths. Slot assignment lives in emit_diagram's pre-pass.
    offset_emu = int(mid_offset_svg * SCALE * 914400)

    if horizontal_first:
        mx_emu = (x1_emu + x2_emu) // 2 + offset_emu
        my_emu = y2_emu
        h_dir = 1 if x1_emu < mx_emu else -1
        v_dir = 1 if y1_emu < y2_emu else -1
        h_dir_2 = 1 if mx_emu < x2_emu else -1
        in_1 = "east" if h_dir == 1 else "west"
        out_1 = "south" if v_dir == 1 else "north"
        in_2 = "south" if v_dir == 1 else "north"
        out_2 = "east" if h_dir_2 == 1 else "west"
    else:
        my_emu = (y1_emu + y2_emu) // 2 + offset_emu
        mx_emu = x2_emu
        v_dir = 1 if y1_emu < my_emu else -1
        h_dir = 1 if x1_emu < x2_emu else -1
        v_dir_2 = 1 if my_emu < y2_emu else -1
        in_1 = "south" if v_dir == 1 else "north"
        out_1 = "east" if h_dir == 1 else "west"
        in_2 = "east" if h_dir == 1 else "west"
        out_2 = "south" if v_dir_2 == 1 else "north"

    stAng_1_deg, swAng_1_deg = ARC_PARAMS[(in_1, out_1)]
    stAng_2_deg, swAng_2_deg = ARC_PARAMS[(in_2, out_2)]
    stAng_1 = stAng_1_deg * 60000
    swAng_1 = swAng_1_deg * 60000
    stAng_2 = stAng_2_deg * 60000
    swAng_2 = swAng_2_deg * 60000

    # Bounding box covering all path points
    min_x = min(x1_emu, x2_emu, mx_emu)
    max_x = max(x1_emu, x2_emu, mx_emu)
    min_y = min(y1_emu, y2_emu, my_emu)
    max_y = max(y1_emu, y2_emu, my_emu)
    box_w = max(1, max_x - min_x)
    box_h = max(1, max_y - min_y)

    # Local coords (offset within box)
    sx, sy = x1_emu - min_x, y1_emu - min_y
    ex, ey = x2_emu - min_x, y2_emu - min_y
    cx_local = mx_emu - min_x
    cy_local = my_emu - min_y

    # Stop-short and post-arc path points
    if horizontal_first:
        p1_x, p1_y = cx_local - r_emu * h_dir, sy           # short of corner 1
        p2_x, p2_y = cx_local, ey - r_emu * v_dir           # short of corner 2
    else:
        p1_x, p1_y = sx, cy_local - r_emu * v_dir
        p2_x, p2_y = ex - r_emu * h_dir, cy_local

    color_hex = _rgb_hex(line_color_rgb)
    dash_xml = f'<a:prstDash val="{dash_val}"/>' if dash_val else '<a:prstDash val="solid"/>'

    _CUSTOM_SHAPE_ID[0] += 1
    shape_id = _CUSTOM_SHAPE_ID[0]

    sp_xml = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr>'
        f'<p:cNvPr id="{shape_id}" name="rounded-z-{shape_id}"/>'
        f'<p:cNvSpPr/><p:nvPr/>'
        f'</p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{min_x}" y="{min_y}"/><a:ext cx="{box_w}" cy="{box_h}"/></a:xfrm>'
        f'<a:custGeom>'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="0" t="0" r="0" b="0"/>'
        f'<a:pathLst>'
        f'<a:path w="{box_w}" h="{box_h}" fill="none">'
        f'<a:moveTo><a:pt x="{sx}" y="{sy}"/></a:moveTo>'
        f'<a:lnTo><a:pt x="{p1_x}" y="{p1_y}"/></a:lnTo>'
        f'<a:arcTo wR="{r_emu}" hR="{r_emu}" stAng="{stAng_1}" swAng="{swAng_1}"/>'
        f'<a:lnTo><a:pt x="{p2_x}" y="{p2_y}"/></a:lnTo>'
        f'<a:arcTo wR="{r_emu}" hR="{r_emu}" stAng="{stAng_2}" swAng="{swAng_2}"/>'
        f'<a:lnTo><a:pt x="{ex}" y="{ey}"/></a:lnTo>'
        f'</a:path>'
        f'</a:pathLst>'
        f'</a:custGeom>'
        f'<a:noFill/>'
        f'<a:ln w="{line_w_emu}" cap="flat">'
        f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
        f'{dash_xml}'
        f'<a:round/>'
        f'<a:tailEnd type="triangle" w="med" len="med"/>'
        f'</a:ln>'
        f'<a:effectLst/>'
        f'</p:spPr>'
        f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
        f'</p:sp>'
    )

    sp_elem = etree.fromstring(sp_xml)
    slide.shapes._spTree.append(sp_elem)


def _new_slide(prs, bg_color="page_bg"):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS[bg_color]
    return slide


# ============================================================================
# Shared chrome — title, lane, caption (used by diagrams + PL slides + decks)
# ============================================================================

def add_title(slide, title, draft_marker=None):
    title_box = slide.shapes.add_textbox(
        Inches(LEFT_MARGIN_IN), Inches(0.2),
        Inches(SLIDE_W_IN - 0.9), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.text = title
    p = tf.paragraphs[0]
    p.font.name = "Lora"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["sand_900"]

    if draft_marker:
        marker_box = slide.shapes.add_textbox(
            Inches(LEFT_MARGIN_IN), Inches(0.55),
            Inches(SLIDE_W_IN - 0.9), Inches(0.3)
        )
        tf2 = marker_box.text_frame
        tf2.margin_left = Inches(0)
        tf2.margin_top = Inches(0)
        tf2.text = draft_marker
        p2 = tf2.paragraphs[0]
        p2.font.name = "Source Code Pro"
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = COLORS["amber_700"]


def add_lane_label(slide, label, y_svg):
    lbl = slide.shapes.add_textbox(
        Inches(LANE_LABEL_X_IN), y_in(y_svg - 18),
        Inches(0.75), Inches(0.25)
    )
    tf = lbl.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.text = label
    p = tf.paragraphs[0]
    p.font.name = "Source Code Pro"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLORS["sand_700"]


def add_lane_divider(slide, y_svg):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(LEFT_MARGIN_IN + 0.5), y_in(y_svg - 5),
        Inches(SLIDE_W_IN - 0.4), y_in(y_svg - 5),
    )
    line.line.color.rgb = COLORS["sand_200"]
    line.line.width = Pt(0.5)
    line.shadow.inherit = False


def add_caption(slide, text):
    box = slide.shapes.add_textbox(
        Inches(LEFT_MARGIN_IN), Inches(SLIDE_H_IN - 0.65),
        Inches(SLIDE_W_IN - 0.9), Inches(0.55)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Lora"
    p.font.size = Pt(9.5)
    p.font.italic = True
    p.font.color.rgb = COLORS["sand_700"]


# ============================================================================
# Diagram-family primitives
# ============================================================================

def add_node(slide, node):
    """Rounded-rectangle node with title, detail lines, and optional watcher."""
    style = MODIFIER_STYLES[node.get("modifier", "default")]
    x, y = x_in(node["x"]), y_in(node["y"])
    w, h = w_in(node["w"]), h_in(node["h"])

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[style["fill"]]
    shape.line.color.rgb = COLORS[style["border"]]
    shape.line.width = Pt(style["border_width"])
    shape.adjustments[0] = 0.06
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = node["title"]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Source Sans 3"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS["sand_900"]

    for line in node.get("detail", []):
        p = tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Source Sans 3"
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS["sand_900"]

    if node.get("watcher"):
        p = tf.add_paragraph()
        p.text = f"▲ WATCHES: {node['watcher']}"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Source Code Pro"
        p.font.size = Pt(6.5)
        p.font.color.rgb = COLORS["sand_700"]

    return shape


def add_callout(slide, node, slide_bboxes=None):
    """Overhead callout to a node, with above→below fallback chain.

    Candidate positions in order:
      1. Above the node (default; bottom-anchored so multi-line stacks up)
      2. Below the node (top-anchored so multi-line stacks down)

    If slide_bboxes provided, picks the first non-colliding candidate; if
    both collide, places above (default) and emits a structured warning.
    Updates slide_bboxes with the chosen bbox so downstream primitives can
    check against it.
    """
    callout = node.get("callout")
    if not callout:
        return
    callout_h_svg = CALLOUT_BOX_H_IN / SCALE   # ~65 svg-px
    gap_svg = CALLOUT_GAP_TO_NODE_IN / SCALE   # ~12 svg-px

    candidates = [
        {
            "bbox": (node["x"], node["y"] - gap_svg - callout_h_svg,
                     node["w"], callout_h_svg),
            "anchor": MSO_ANCHOR.BOTTOM,
            "label": "above",
        },
        {
            "bbox": (node["x"], node["y"] + node["h"] + gap_svg,
                     node["w"], callout_h_svg),
            "anchor": MSO_ANCHOR.TOP,
            "label": "below",
        },
    ]

    chosen = candidates[0]
    if slide_bboxes is not None:
        for cand in candidates:
            if not _collides_any(cand["bbox"], slide_bboxes):
                chosen = cand
                break
        else:
            print(f"  ⚠ callout for node '{node.get('id', '?')}' collides "
                  f"in all candidates — placed above (default)", file=sys.stderr)

    cx, cy, cw, ch = chosen["bbox"]
    box = slide.shapes.add_textbox(x_in(cx), y_in(cy), w_in(cw), h_in(ch))
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = chosen["anchor"]
    tf.word_wrap = True
    tf.text = callout["text"]
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Source Code Pro"
    p.font.size = Pt(7.5)
    p.font.bold = True
    p.font.color.rgb = COLORS[CALLOUT_COLORS.get(callout["type"], "amber_700")]

    if slide_bboxes is not None:
        slide_bboxes.append(chosen["bbox"])


def add_edge(slide, edge, nodes_by_id, slot=0, slide_bboxes=None):
    """Connector between two nodes — dominant-axis routing (v0.3).

    Pick exit/entry sides by the dominant displacement axis:
    - |Δx| ≥ |Δy| (horizontal-dominant): exit right or left, enter left or
      right. Continues horizontal flow; fork patterns (one source, multiple
      destinations stacked vertically to the side) all exit from the same
      side and route via Z-shape elbows.
    - |Δy| > |Δx| (vertical-dominant): exit bottom or top, enter top or
      bottom. Continues vertical flow.

    Connector implementation: STRAIGHT when the secondary-axis offset is below
    LINEAR_TOLERANCE_SVG (essentially-linear path); custom-geometry rounded Z
    shape otherwise (replaces MSO_CONNECTOR_TYPE.ELBOW because the preset's
    bend orientation is determined by its bounding box, not by which shape
    edges the endpoints touch — so it can't be forced to enter the destination
    perpendicular to a side edge). The rounded Z gives both perpendicular
    entry AND rounded corners (CORNER_RADIUS_PT).

    All edges carry triangle arrow head + suppressed shadow. Muted edges
    render dashed via prstDash. Labels mask the line via page_bg fill sized
    to the text width.
    """
    src = nodes_by_id[edge["from"]]
    dst = nodes_by_id[edge["to"]]

    src_cx = src["x"] + src["w"] / 2
    src_cy = src["y"] + src["h"] / 2
    dst_cx = dst["x"] + dst["w"] / 2
    dst_cy = dst["y"] + dst["h"] / 2

    dx = dst_cx - src_cx
    dy = dst_cy - src_cy
    horizontal_dominant = abs(dx) >= abs(dy)

    if horizontal_dominant:
        if dx >= 0:
            x1, y1 = src["x"] + src["w"], src_cy
            x2, y2 = dst["x"], dst_cy
        else:
            x1, y1 = src["x"], src_cy
            x2, y2 = dst["x"] + dst["w"], dst_cy
        secondary_offset = abs(y2 - y1)
    else:
        if dy >= 0:
            x1, y1 = src_cx, src["y"] + src["h"]
            x2, y2 = dst_cx, dst["y"]
        else:
            x1, y1 = src_cx, src["y"]
            x2, y2 = dst_cx, dst["y"] + dst["h"]
        secondary_offset = abs(x2 - x1)

    style = edge.get("style", "default")
    if style == "emphasis":
        color = COLORS["teal_700"]
        width_pt = 1.75
        dash_val = None
    elif style == "muted":
        color = COLORS["arrow_muted"]
        width_pt = 1.25
        dash_val = "sysDash"
    else:
        color = COLORS["arrow_default"]
        width_pt = 1.5
        dash_val = None

    if secondary_offset < LINEAR_TOLERANCE_SVG:
        # Essentially linear — use stock STRAIGHT connector
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            x_in(x1), y_in(y1), x_in(x2), y_in(y2),
        )
        conn.shadow.inherit = False
        conn.line.color.rgb = color
        conn.line.width = Pt(width_pt)
        if dash_val:
            _set_line_dash(conn.line, dash_val)
        _add_arrow_tail(conn.line, kind="triangle", width="med", length="med")
    else:
        # Non-linear — use custom-geometry rounded Z. horizontal_first matches
        # the dominant axis so the last segment enters the destination edge
        # perpendicular (not parallel).
        _add_rounded_z_connector(
            slide, x1, y1, x2, y2,
            line_color_rgb=color,
            line_width_pt=width_pt,
            dash_val=dash_val,
            horizontal_first=horizontal_dominant,
            mid_offset_svg=slot * OVERLAP_OFFSET_SVG,
        )

    if edge.get("label"):
        _add_edge_label_masked(slide, x1, y1, x2, y2, edge["label"],
                               slide_bboxes=slide_bboxes)


def _add_edge_label_masked(slide, x1_svg, y1_svg, x2_svg, y2_svg, text, slide_bboxes=None):
    """Edge label with page_bg fill, sized to text width + small padding.

    Fallback chain of candidate positions along/near the edge:
      1. Geometric midpoint
      2. 1/3 along path (closer to source)
      3. 2/3 along path (closer to destination)
      4. Perpendicular offset above the path (~35 svg-px)
      5. Perpendicular offset below the path (~35 svg-px)

    If slide_bboxes provided, picks the first non-colliding candidate; if
    all collide, places at midpoint (default) and emits a structured warning.
    """
    char_w_in = 0.058
    h_pad_in = 0.08
    label_w_in = max(0.25, len(text) * char_w_in + h_pad_in)
    label_h_in = 0.18
    label_w_svg = label_w_in / SCALE
    label_h_svg = label_h_in / SCALE

    mx = (x1_svg + x2_svg) / 2
    my = (y1_svg + y2_svg) / 2
    dx, dy = x2_svg - x1_svg, y2_svg - y1_svg
    length = max(1.0, (dx * dx + dy * dy) ** 0.5)
    # Unit perpendicular vector
    pux = -dy / length
    puy = dx / length

    # Candidate set: along-path positions for in-gap cases, then perpendicular
    # sweep at growing offsets for cases where the gap is too narrow for the
    # label. The sweep escapes the node bbox by ~70-100 svg-px (~0.6-0.85in),
    # enough to clear most lane content vertically or horizontally.
    candidate_centers = [
        ("midpoint", (mx, my)),
        ("1/3-src", (x1_svg + dx / 3, y1_svg + dy / 3)),
        ("2/3-dst", (x1_svg + 2 * dx / 3, y1_svg + 2 * dy / 3)),
    ]
    for dist in (45, 75, 105):
        candidate_centers.append((f"perp+{dist}", (mx + pux * dist, my + puy * dist)))
        candidate_centers.append((f"perp-{dist}", (mx - pux * dist, my - puy * dist)))

    chosen_label, chosen_xy = candidate_centers[0]
    if slide_bboxes is not None:
        for label_name, (cx, cy) in candidate_centers:
            cand_bbox = (cx - label_w_svg / 2, cy - label_h_svg / 2,
                         label_w_svg, label_h_svg)
            if not _collides_any(cand_bbox, slide_bboxes):
                chosen_label, chosen_xy = label_name, (cx, cy)
                break
        else:
            print(f"  ⚠ edge label '{text}' collides in all candidates "
                  f"— placed at midpoint", file=sys.stderr)

    cx, cy = chosen_xy
    box = slide.shapes.add_textbox(
        x_in(cx) - Inches(label_w_in / 2),
        y_in(cy) - Inches(label_h_in / 2),
        Inches(label_w_in),
        Inches(label_h_in),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["page_bg"]
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Source Code Pro"
    p.font.size = Pt(7)
    p.font.color.rgb = COLORS["sand_700"]

    if slide_bboxes is not None:
        slide_bboxes.append((cx - label_w_svg / 2, cy - label_h_svg / 2,
                              label_w_svg, label_h_svg))


# ============================================================================
# Workflow diagram emitter (single-slide)
# ============================================================================

def apply_grid_layout(spec, canvas_w_svg=1450, h_pad_svg=80, cluster_tol_svg=40):
    """Recompute node x_center as a column grid spanning the full slide canvas.

    Clusters source x_center values into columns (nodes within cluster_tol_svg
    of each other are the same column), then redistributes column centers
    evenly across the usable canvas width. Source intent (which-column-which-
    lane, relative ordering) is preserved; absolute coordinates are recomputed
    to give every column equal horizontal width budget.

    Why: source x_centers were authored for a notional 1100-svg viewBox with
    ~220 svg-px between columns — too tight for any meaningful edge label.
    The slide canvas is wider (~1450 svg usable); equal redistribution buys
    each gap enough room for ~17-33-char labels depending on column count.
    """
    nodes = spec.get("nodes", [])
    if not nodes:
        return spec

    def cluster_key(n):
        return round((n["x"] + n["w"] / 2) / cluster_tol_svg) * cluster_tol_svg

    col_keys = sorted({cluster_key(n) for n in nodes})
    col_id = {k: i for i, k in enumerate(col_keys)}
    num_cols = len(col_keys)

    avail = canvas_w_svg - 2 * h_pad_svg
    if num_cols == 1:
        targets = [canvas_w_svg / 2]
    else:
        spacing = avail / (num_cols - 1)
        targets = [h_pad_svg + spacing * i for i in range(num_cols)]

    for n in nodes:
        n["x"] = targets[col_id[cluster_key(n)]] - n["w"] / 2

    return spec


def render_workflow_to_slide(slide, spec):
    """Render one workflow spec onto one slide.

    Maintains a slide_bboxes ledger of placed-primitive bounding boxes (in
    svg-px). Each primitive declares a fallback chain of candidate positions
    and the renderer picks the first non-colliding one. Unavoidable collisions
    emit a structured warning to stderr; the slide still renders.

    Layout pre-pass: apply_grid_layout redistributes source x_centers across
    the canvas so every column gets equal width budget.
    """
    apply_grid_layout(spec)
    add_title(slide, spec["title"], spec.get("draft_marker"))

    for lane in spec["lanes"]:
        add_lane_label(slide, lane["label"], lane["y"])
        add_lane_divider(slide, lane["y"])

    nodes_by_id = {n["id"]: n for n in spec["nodes"]}

    # Pre-pass: assign slot offsets to edges sharing an undirected node pair,
    # so opposing-direction or parallel edges don't trace identical Z paths.
    pair_counts = {}
    for edge in spec["edges"]:
        key = frozenset((edge["from"], edge["to"]))
        pair_counts[key] = pair_counts.get(key, 0) + 1
    pair_seen = {}

    # Collision ledger: nodes go in first as immutable anchors, then callouts
    # with above/below fallback, then edge labels with along-path / perpendicular
    # fallback. Each primitive checks against everything already placed.
    slide_bboxes = []
    for node in spec["nodes"]:
        add_node(slide, node)
        slide_bboxes.append(_node_bbox(node))
    for node in spec["nodes"]:
        add_callout(slide, node, slide_bboxes=slide_bboxes)
    for edge in spec["edges"]:
        key = frozenset((edge["from"], edge["to"]))
        n = pair_counts[key]
        i = pair_seen.get(key, 0)
        slot = (i - (n - 1) / 2) * 2 if n > 1 else 0
        pair_seen[key] = i + 1
        add_edge(slide, edge, nodes_by_id, slot=slot, slide_bboxes=slide_bboxes)

    add_caption(slide, spec.get("caption", ""))


def emit_diagram(spec, output_path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = _new_slide(prs)
    render_workflow_to_slide(slide, spec)
    prs.save(output_path)
    print(f"  ✓ {output_path}")


def emit_deck(specs, output_path, deck_title=None, deck_subtitle=None):
    """Render a multi-slide deck — one slide per workflow spec.

    Optionally inserts a cover slide (deck-title-stub treatment) at the front
    when deck_title is provided.
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    if deck_title:
        cover = _new_slide(prs)
        title_box = cover.shapes.add_textbox(
            Inches(LEFT_MARGIN_IN), Inches(2.4),
            Inches(SLIDE_W_IN - 0.9), Inches(1.0)
        )
        tf = title_box.text_frame
        tf.margin_left = Inches(0)
        tf.text = deck_title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Lora"
        p.font.size = Pt(32)
        p.font.color.rgb = COLORS["sand_900"]
        if deck_subtitle:
            sub_box = cover.shapes.add_textbox(
                Inches(LEFT_MARGIN_IN), Inches(3.5),
                Inches(SLIDE_W_IN - 0.9), Inches(0.5)
            )
            tf2 = sub_box.text_frame
            tf2.margin_left = Inches(0)
            tf2.text = deck_subtitle
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            p2.font.name = "Lora"
            p2.font.size = Pt(13)
            p2.font.italic = True
            p2.font.color.rgb = COLORS["sand_700"]

    for spec in specs:
        slide = _new_slide(prs)
        render_workflow_to_slide(slide, spec)

    prs.save(output_path)
    print(f"  ✓ {output_path} ({len(prs.slides)} slides)")


# ============================================================================
# Slides Pattern Library — shared eyebrow/head and section header
# ============================================================================

def _pl_eyebrow(slide, eyebrow):
    box = slide.shapes.add_textbox(
        Inches(LEFT_MARGIN_IN), Inches(0.3),
        Inches(SLIDE_W_IN - 0.9), Inches(0.3)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.text = eyebrow
    p = tf.paragraphs[0]
    p.font.name = "Source Code Pro"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLORS["sand_700"]


def _pl_head(slide, head):
    box = slide.shapes.add_textbox(
        Inches(LEFT_MARGIN_IN), Inches(0.6),
        Inches(SLIDE_W_IN - 0.9), Inches(0.6)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.text = head
    p = tf.paragraphs[0]
    p.font.name = "Lora"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS["sand_900"]


# ============================================================================
# PL: Cover
# ============================================================================

def pl_slide_cover(prs):
    slide = _new_slide(prs)
    title_box = slide.shapes.add_textbox(
        Inches(LEFT_MARGIN_IN), Inches(2.4),
        Inches(SLIDE_W_IN - 0.9), Inches(1.0)
    )
    tf = title_box.text_frame
    tf.margin_left = Inches(0)
    tf.text = "Haven — Slides Pattern Library"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Lora"
    p.font.size = Pt(34)
    p.font.color.rgb = COLORS["sand_900"]

    sub_box = slide.shapes.add_textbox(
        Inches(LEFT_MARGIN_IN), Inches(3.5),
        Inches(SLIDE_W_IN - 0.9), Inches(0.5)
    )
    tf2 = sub_box.text_frame
    tf2.margin_left = Inches(0)
    tf2.text = "Canonical visual treatments — diagram primitives + deck primitives (stubs)"
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Lora"
    p2.font.size = Pt(13)
    p2.font.italic = True
    p2.font.color.rgb = COLORS["sand_700"]

    marker_box = slide.shapes.add_textbox(
        Inches(LEFT_MARGIN_IN), Inches(4.2),
        Inches(SLIDE_W_IN - 0.9), Inches(0.4)
    )
    tf3 = marker_box.text_frame
    tf3.margin_left = Inches(0)
    tf3.text = "DRAFT • v0.2 • spec source for all Haven .pptx output"
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = "Source Code Pro"
    p3.font.size = Pt(10)
    p3.font.color.rgb = COLORS["amber_700"]


# ============================================================================
# PL: Diagram-family — comparison slides (replace per-primitive variant slides)
# ============================================================================

def _canvas_w_svg():
    """SVG-px equivalent of the horizontal canvas (between left/right margins)."""
    return (SLIDE_W_IN - 2 * LEFT_MARGIN_IN) / SCALE


def pl_slide_modifier_comparison(prs):
    """All 5 node modifiers side-by-side with their names above each."""
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DIAGRAM — NODE MODIFIERS")
    _pl_head(slide, "Five borders, scaled by attention weight")

    modifiers = [
        ("default",              "Confirm meal order",    "selection captured"),
        ("uncertain-assumption", "Hand-off",              "BHN role assumed"),
        ("uncertain-tbd",        "Threshold trigger",     "exact threshold pending"),
        ("uncertain-gap",        "Severity assessment",   "tier boundaries undefined"),
        ("attestation-gate",     "Attestation gate",      "clinical lead signs off"),
    ]

    node_w_svg = 220
    node_h_svg = 110
    gap_svg = 30
    total_w = 5 * node_w_svg + 4 * gap_svg
    start_x = (_canvas_w_svg() - total_w) / 2
    name_y = 280
    node_y = 320

    for i, (mod, title, detail) in enumerate(modifiers):
        x = start_x + i * (node_w_svg + gap_svg)
        # Modifier name above
        name_box = slide.shapes.add_textbox(
            x_in(x), y_in(name_y),
            w_in(node_w_svg), Inches(0.25)
        )
        tf = name_box.text_frame
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        tf.text = mod
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Source Code Pro"
        p.font.size = Pt(8.5)
        p.font.bold = True
        p.font.color.rgb = COLORS["sand_700"]
        # Node
        add_node(slide, {
            "id": f"mod-{i}",
            "x": x, "y": node_y, "w": node_w_svg, "h": node_h_svg,
            "title": title, "detail": [detail], "modifier": mod,
        })

    add_caption(slide,
        "Default and uncertain-assumption are quiet (sand borders). uncertain-tbd carries amber attention. uncertain-gap carries red weight for structurally-required absences. attestation-gate gets teal heavy border — reserved for accountable-human sign-off. Pair each with the matching callout on the next slide.")


def pl_slide_callout_comparison(prs):
    """All 4 callout types over default-modifier nodes."""
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DIAGRAM — CALLOUTS")
    _pl_head(slide, "Four types, scaled by attention weight")

    callouts = [
        ("assumption", "[ASSUMPTION] BHN role unresolved"),
        ("tbd",        "[TBD] threshold pending HDG"),
        ("gap",        "[GAP] tier boundaries pending Marrero"),
        ("critical",   "[TBD CRITICAL] first action — Vanessa/Marrero"),
    ]

    node_w_svg = 260
    node_h_svg = 80
    gap_svg = 40
    total_w = 4 * node_w_svg + 3 * gap_svg
    start_x = (_canvas_w_svg() - total_w) / 2
    type_name_y = 200
    node_y = 360

    for i, (callout_type, callout_text) in enumerate(callouts):
        x = start_x + i * (node_w_svg + gap_svg)
        # Type name at top of the column
        name_box = slide.shapes.add_textbox(
            x_in(x), y_in(type_name_y),
            w_in(node_w_svg), Inches(0.25)
        )
        tf = name_box.text_frame
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        tf.text = callout_type
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Source Code Pro"
        p.font.size = Pt(8.5)
        p.font.bold = True
        p.font.color.rgb = COLORS["sand_700"]
        # Node + callout
        node = {
            "id": f"call-{i}",
            "x": x, "y": node_y, "w": node_w_svg, "h": node_h_svg,
            "title": "Example step", "detail": ["pairs with callout above"],
            "modifier": "default",
            "callout": {"text": callout_text, "type": callout_type},
        }
        add_node(slide, node)
        add_callout(slide, node)

    add_caption(slide,
        "Quietest → loudest: assumption (sand-700) < tbd (amber-700) < gap (red-500) < critical (red-700, bold). Pair with the matching node modifier. The callout names what is pending; the modifier signals weight on the step itself.")


def pl_slide_edge_comparison(prs):
    """All 3 edge styles, each a labeled row with src→dst demo."""
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DIAGRAM — EDGE STYLES")
    _pl_head(slide, "Three weights for flow")

    styles = [
        ("default",  "Sequential flow. Warm-amber arrow, 1.5pt, triangle head."),
        ("emphasis", "Primary / load-bearing path. Teal-700, 1.75pt."),
        ("muted",    "Background or fallback. Lighter warm-amber, 1.25pt dashed."),
    ]

    row_h_svg = 110
    row_start_y = 220
    label_x = 60
    label_w = 260
    demo_src_x = 380
    demo_dst_x = 880
    demo_node_w = 180
    demo_node_h = 60

    for i, (style, descr) in enumerate(styles):
        row_y = row_start_y + i * row_h_svg

        # Left: style name + description (Lora + SS3)
        label_box = slide.shapes.add_textbox(
            x_in(label_x), y_in(row_y),
            w_in(label_w), h_in(demo_node_h + 30)
        )
        tf = label_box.text_frame
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.text = style
        p = tf.paragraphs[0]
        p.font.name = "Lora"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["sand_900"]
        p2 = tf.add_paragraph()
        p2.text = descr
        p2.font.name = "Source Sans 3"
        p2.font.size = Pt(9)
        p2.font.color.rgb = COLORS["sand_700"]

        # Right: src → dst demo
        src = {
            "id": f"e{i}s", "x": demo_src_x, "y": row_y,
            "w": demo_node_w, "h": demo_node_h,
            "title": "A", "detail": [], "modifier": "default",
        }
        dst = {
            "id": f"e{i}d", "x": demo_dst_x, "y": row_y,
            "w": demo_node_w, "h": demo_node_h,
            "title": "B", "detail": [], "modifier": "default",
        }
        add_node(slide, src)
        add_node(slide, dst)
        add_edge(slide,
                 {"from": f"e{i}s", "to": f"e{i}d", "style": style},
                 {f"e{i}s": src, f"e{i}d": dst})

    add_caption(slide,
        "Default reads as standard step-to-step succession. Emphasis lifts the critical-path edge above default. Muted recedes into background — timeouts, deferred handlers, lower-traffic branches. Dash on muted is sysDash (Slides does not expose dash spacing).")


# ============================================================================
# PL: Diagram-family — page chrome (kept one-per-primitive; genuinely different)
# ============================================================================

def pl_slide_chrome_lane_label(prs):
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DIAGRAM — PAGE CHROME")
    _pl_head(slide, "Lane label + divider")
    add_lane_label(slide, "AGENT", 220)
    add_lane_divider(slide, 220)
    add_node(slide, {
        "id": "a", "x": 220, "y": 200, "w": 220, "h": 80,
        "title": "Example agent step", "detail": ["lane = AGENT"],
        "modifier": "default",
    })
    add_lane_label(slide, "HUMAN", 400)
    add_lane_divider(slide, 400)
    add_node(slide, {
        "id": "b", "x": 220, "y": 380, "w": 220, "h": 80,
        "title": "Example human step", "detail": ["lane = HUMAN"],
        "modifier": "default",
    })
    add_caption(slide,
        "Lane labels run in Source Code Pro 8pt bold (sand-700) at the page's left margin. A 0.5pt sand-200 divider extends from after the label across the canvas. Lane roles map 1:1 to participating actor classes (AGENT / HUMAN / PARTNER / SYSTEM).")


def pl_slide_chrome_watcher_label(prs):
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DIAGRAM — PAGE CHROME")
    _pl_head(slide, "Watcher label (▲ WATCHES:)")
    add_node(slide, {
        "id": "demo", "x": 460, "y": 305, "w": 260, "h": 120,
        "title": "Confirm meal order",
        "detail": ["selection captured", "queue → kitchen"],
        "watcher": "clinical-lead",
        "modifier": "default",
    })
    add_caption(slide,
        "Watcher labels run as the last line inside the node's text frame in Source Code Pro 6.5pt sand-700, prefixed with ▲. They name an accountable role that monitors the step without owning it. Distinct from attestation-gate, which owns sign-off rather than watching.")


def pl_slide_chrome_title(prs):
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DIAGRAM — PAGE CHROME")
    _pl_head(slide, "Title + draft marker")
    demo_title_box = slide.shapes.add_textbox(
        Inches(LEFT_MARGIN_IN + 1.5), Inches(2.8),
        Inches(SLIDE_W_IN - 4.0), Inches(0.5)
    )
    tf = demo_title_box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.text = "Workflow — example title"
    p = tf.paragraphs[0]
    p.font.name = "Lora"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["sand_900"]
    demo_marker_box = slide.shapes.add_textbox(
        Inches(LEFT_MARGIN_IN + 1.5), Inches(3.35),
        Inches(SLIDE_W_IN - 4.0), Inches(0.3)
    )
    tf2 = demo_marker_box.text_frame
    tf2.margin_left = Inches(0)
    tf2.margin_top = Inches(0)
    tf2.text = "DRAFT • v0.1 • pending Marrero + Healthcare Data Governance review"
    p2 = tf2.paragraphs[0]
    p2.font.name = "Source Code Pro"
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = COLORS["amber_700"]
    add_caption(slide,
        "Title in Lora 20pt sand-900 at top-left. Draft marker beneath in Source Code Pro 8.5pt amber-700: version + draft state + named reviewers. Amber signals draft status without obscuring the title.")


def pl_slide_chrome_caption(prs):
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DIAGRAM — PAGE CHROME")
    _pl_head(slide, "Caption (bottom)")
    add_node(slide, {
        "id": "demo", "x": 460, "y": 280, "w": 260, "h": 80,
        "title": "Diagram content",
        "detail": ["nodes and edges live above the caption"],
        "modifier": "default",
    })
    add_caption(slide,
        "Caption sits at slide bottom in Lora 9.5pt italic sand-700, centered. One or two short sentences summarizing the path or naming the emphasized branch. Reserve for the reading thesis of the diagram, not new content.")


# ============================================================================
# PL: Diagram-family — compositions
# ============================================================================

def pl_slide_composition_title_only(prs):
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DIAGRAM — COMPOSITION")
    _pl_head(slide, "Title-only diagram")
    add_node(slide, {
        "id": "demo", "x": 460, "y": 305, "w": 260, "h": 120,
        "title": "Step", "detail": ["sparse compositions still valid"],
        "modifier": "default",
    })
    add_caption(slide,
        "Smallest valid composition: eyebrow + head + at least one node + caption. Lane labels optional when one actor class participates. Restraint: do not pad sparse diagrams with chrome they do not need.")


def pl_slide_composition_single_node_callout(prs):
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DIAGRAM — COMPOSITION")
    _pl_head(slide, "Single node + callout")
    node = {
        "id": "demo", "x": 460, "y": 320, "w": 260, "h": 120,
        "title": "Confirm escalation tier",
        "detail": ["red / yellow / green decision", "based on PHQ-9 severity"],
        "modifier": "uncertain-tbd",
        "callout": {"text": "[TBD] tier boundaries pending Marrero", "type": "tbd"},
    }
    add_node(slide, node)
    add_callout(slide, node)
    add_caption(slide,
        "The atomic unit of decorated process. Node carries the work; callout names what is pending or unresolved about it. Callout anchored bottom-up: 0.10in gap to node top, settles consistent regardless of wrap.")


def pl_slide_composition_cross_lane_edge(prs):
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DIAGRAM — COMPOSITION")
    _pl_head(slide, "Cross-lane edge with label")
    add_lane_label(slide, "AGENT", 220)
    add_lane_divider(slide, 220)
    add_lane_label(slide, "HUMAN", 440)
    add_lane_divider(slide, 440)
    src = {
        "id": "a", "x": 280, "y": 200, "w": 240, "h": 90,
        "title": "Agent step", "detail": ["produces decision draft"],
        "modifier": "default",
    }
    dst = {
        "id": "b", "x": 700, "y": 420, "w": 240, "h": 90,
        "title": "Human attestation", "detail": ["clinical lead signs off"],
        "modifier": "attestation-gate",
    }
    add_node(slide, src)
    add_node(slide, dst)
    add_edge(slide,
             {"from": "a", "to": "b", "label": "draft → attestation", "style": "emphasis"},
             {"a": src, "b": dst})
    add_caption(slide,
        "Cross-lane edges route by dominant axis. Horizontal-dominant (|Δx| ≥ |Δy|): exit right of source, enter left of destination via Z-shape ELBOW. Vertical-dominant: exit bottom, enter top. STRAIGHT when secondary-axis offset is near-zero. Triangle arrow head at destination; label mask sized to text width.")


# ============================================================================
# PL: Deck-family stubs (stake the convergence shape, B11 builds out)
# ============================================================================

def pl_slide_deck_title_stub(prs):
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DECK — STUB")
    _pl_head(slide, "Title slide treatment")

    # Demonstrate a deck title-slide-in-miniature inside the canvas
    frame_l = Inches(2.5)
    frame_t = Inches(2.1)
    frame_w = Inches(SLIDE_W_IN - 5.0)
    frame_h = Inches(3.2)

    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, frame_l, frame_t, frame_w, frame_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = COLORS["page_bg"]
    frame.line.color.rgb = COLORS["sand_300"]
    frame.line.width = Pt(0.5)
    frame.shadow.inherit = False

    # Inner content
    inner_l = frame_l + Inches(0.5)
    inner_w = frame_w - Inches(1.0)

    title_box = slide.shapes.add_textbox(inner_l, frame_t + Inches(0.8), inner_w, Inches(0.7))
    tf = title_box.text_frame
    tf.margin_left = Inches(0); tf.margin_top = Inches(0)
    tf.text = "Workflow Drilldowns — Cena Health"
    p = tf.paragraphs[0]
    p.font.name = "Lora"; p.font.size = Pt(22)
    p.font.color.rgb = COLORS["sand_900"]

    sub_box = slide.shapes.add_textbox(inner_l, frame_t + Inches(1.7), inner_w, Inches(0.5))
    tf2 = sub_box.text_frame
    tf2.margin_left = Inches(0); tf2.margin_top = Inches(0)
    tf2.text = "Patient-app v0.1 — early-clinical-pilot review"
    p2 = tf2.paragraphs[0]
    p2.font.name = "Lora"; p2.font.size = Pt(12)
    p2.font.italic = True
    p2.font.color.rgb = COLORS["sand_700"]

    marker_box = slide.shapes.add_textbox(
        inner_l, frame_t + frame_h - Inches(0.6), inner_w, Inches(0.3))
    tf3 = marker_box.text_frame
    tf3.margin_left = Inches(0); tf3.margin_top = Inches(0)
    tf3.text = "DRAFT • 2026-06-04 • Aaron Sleeper"
    p3 = tf3.paragraphs[0]
    p3.font.name = "Source Code Pro"; p3.font.size = Pt(8.5)
    p3.font.color.rgb = COLORS["amber_700"]

    add_caption(slide,
        "Deck title-slide stub. Lora 32-36pt title at full size; Lora italic 13pt subtitle; Source Code Pro 8.5pt draft marker amber-700. Same typography + palette + draft-marker convention as diagram chrome — shared base layer pays off. Full deck primitives land in B11.")


def pl_slide_deck_content_stub(prs):
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DECK — STUB")
    _pl_head(slide, "Content slide (single column + bullets)")

    # Demonstrate a content-slide-in-miniature
    frame_l = Inches(2.5)
    frame_t = Inches(2.0)
    frame_w = Inches(SLIDE_W_IN - 5.0)
    frame_h = Inches(3.4)

    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, frame_l, frame_t, frame_w, frame_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = COLORS["page_bg"]
    frame.line.color.rgb = COLORS["sand_300"]
    frame.line.width = Pt(0.5)
    frame.shadow.inherit = False

    inner_l = frame_l + Inches(0.4)
    inner_w = frame_w - Inches(0.8)

    head_box = slide.shapes.add_textbox(inner_l, frame_t + Inches(0.3), inner_w, Inches(0.6))
    tf = head_box.text_frame
    tf.margin_left = Inches(0); tf.margin_top = Inches(0)
    tf.text = "Escalation path — at a glance"
    p = tf.paragraphs[0]
    p.font.name = "Lora"; p.font.size = Pt(20)
    p.font.color.rgb = COLORS["sand_900"]

    bullets_box = slide.shapes.add_textbox(
        inner_l, frame_t + Inches(1.1), inner_w, frame_h - Inches(1.4))
    tf2 = bullets_box.text_frame
    tf2.margin_left = Inches(0); tf2.margin_top = Inches(0)
    tf2.word_wrap = True
    bullets = [
        "Trigger captures PHQ-9 positive response → severity assessment",
        "Severity routes to clinical-lead attestation when red or yellow",
        "Hand-off to BHN / UConn clinical team after signed attestation",
        "Audit record sealed at every terminal state (chain of custody)",
    ]
    for i, b in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"•  {b}"
        p.font.name = "Source Sans 3"
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS["sand_900"]
        p.space_after = Pt(6)

    add_caption(slide,
        "Deck content-slide stub. Lora 20-24pt heading, Source Sans 3 13pt bullets with breathing-room spacing. The framed box represents the slide canvas at miniature scale; production deck slides use the full canvas with comparable proportions.")


def pl_slide_deck_two_column_stub(prs):
    slide = _new_slide(prs)
    _pl_eyebrow(slide, "DECK — STUB")
    _pl_head(slide, "Two-column (image / copy)")

    frame_l = Inches(2.0)
    frame_t = Inches(2.0)
    frame_w = Inches(SLIDE_W_IN - 4.0)
    frame_h = Inches(3.4)

    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, frame_l, frame_t, frame_w, frame_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = COLORS["page_bg"]
    frame.line.color.rgb = COLORS["sand_300"]
    frame.line.width = Pt(0.5)
    frame.shadow.inherit = False

    inner_l = frame_l + Inches(0.4)
    inner_w = frame_w - Inches(0.8)
    inner_t = frame_t + Inches(0.3)

    # Header spanning both columns
    head_box = slide.shapes.add_textbox(inner_l, inner_t, inner_w, Inches(0.5))
    tf = head_box.text_frame
    tf.margin_left = Inches(0); tf.margin_top = Inches(0)
    tf.text = "Severity assessment — boundaries"
    p = tf.paragraphs[0]
    p.font.name = "Lora"; p.font.size = Pt(18)
    p.font.color.rgb = COLORS["sand_900"]

    # Two columns below header
    col_top = inner_t + Inches(0.85)
    col_h = frame_h - Inches(1.4)
    col_gap = Inches(0.3)
    col_w = (inner_w - col_gap) / 2

    # Left: image placeholder
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, inner_l, col_top, col_w, col_h)
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = COLORS["sand_100"]
    placeholder.line.color.rgb = COLORS["sand_300"]
    placeholder.line.width = Pt(0.5)
    placeholder.shadow.inherit = False
    tf_ph = placeholder.text_frame
    tf_ph.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_ph.text = "image / chart"
    p_ph = tf_ph.paragraphs[0]
    p_ph.alignment = PP_ALIGN.CENTER
    p_ph.font.name = "Source Code Pro"; p_ph.font.size = Pt(9)
    p_ph.font.color.rgb = COLORS["sand_700"]

    # Right: copy
    copy_box = slide.shapes.add_textbox(
        inner_l + col_w + col_gap, col_top, col_w, col_h)
    tf_c = copy_box.text_frame
    tf_c.margin_left = Inches(0); tf_c.margin_top = Inches(0)
    tf_c.word_wrap = True
    tf_c.text = "Tier boundaries are pending Marrero's clinical review."
    p_c = tf_c.paragraphs[0]
    p_c.font.name = "Source Sans 3"; p_c.font.size = Pt(12)
    p_c.font.color.rgb = COLORS["sand_900"]
    p_c2 = tf_c.add_paragraph()
    p_c2.text = ""
    p_c2.space_after = Pt(8)
    p_c3 = tf_c.add_paragraph()
    p_c3.text = "Working assumption: red ≥ 15 routes to immediate attestation; yellow 10-14 routes to within-24h queue."
    p_c3.font.name = "Source Sans 3"; p_c3.font.size = Pt(11)
    p_c3.font.color.rgb = COLORS["sand_700"]

    add_caption(slide,
        "Deck two-column stub. Header spans both columns; image / chart placeholder left, copy right. Equal-width columns with 0.3in gap. Same Lora head + Source Sans 3 body as the content stub; same palette as the diagram family.")


# ============================================================================
# PL emitter
# ============================================================================

def emit_pl(output_path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    pl_slide_cover(prs)

    # Diagram family — variant primitives as comparison
    pl_slide_modifier_comparison(prs)
    pl_slide_callout_comparison(prs)
    pl_slide_edge_comparison(prs)

    # Diagram family — page chrome (kept individual, genuinely different)
    pl_slide_chrome_lane_label(prs)
    pl_slide_chrome_watcher_label(prs)
    pl_slide_chrome_title(prs)
    pl_slide_chrome_caption(prs)

    # Diagram family — compositions
    pl_slide_composition_title_only(prs)
    pl_slide_composition_single_node_callout(prs)
    pl_slide_composition_cross_lane_edge(prs)

    # Deck family — stubs (B11 builds out)
    pl_slide_deck_title_stub(prs)
    pl_slide_deck_content_stub(prs)
    pl_slide_deck_two_column_stub(prs)

    prs.save(output_path)
    print(f"  ✓ {output_path} ({len(prs.slides)} slides)")


# ============================================================================
# Workflow spec (escalation — PHQ-9 positive)
# ============================================================================

ESCALATION_SPEC = {
    "title": "Escalation — PHQ9 positive response",
    "draft_marker": "DRAFT • v0.1 • pending Marrero + Healthcare Data Governance review",
    "lanes": [
        {"label": "AGENT",   "y": 100},
        {"label": "HUMAN",   "y": 220},
        {"label": "PARTNER", "y": 340},
        {"label": "SYSTEM",  "y": 460},
    ],
    "nodes": [
        {
            "id": "trigger",
            "x": 100, "y": 65, "w": 180, "h": 85,
            "title": "Trigger",
            "detail": ["PHQ9 positive", "response detected"],
            "watcher": "clinical-lead",
            "modifier": "uncertain-tbd",
            "callout": {"text": "[TBD] threshold pending HDG", "type": "tbd"},
        },
        {
            "id": "decision",
            "x": 340, "y": 65, "w": 180, "h": 85,
            "title": "Decision",
            "detail": ["Severity assessment", "red / yellow / green"],
            "watcher": "clinical-lead",
            "modifier": "uncertain-gap",
            "callout": {"text": "[GAP] tier boundaries pending Marrero", "type": "gap"},
        },
        {
            "id": "attestation-gate",
            "x": 580, "y": 180, "w": 200, "h": 100,
            "title": "Attestation gate",
            "detail": ["Clinical lead review", "evidence + decision + rationale",
                       "+ SLA timeout fallback"],
            "modifier": "attestation-gate",
            "callout": {"text": "[TBD CRITICAL] first-action specifics — Vanessa/Marrero",
                        "type": "critical"},
        },
        {
            "id": "handoff",
            "x": 860, "y": 295, "w": 200, "h": 90,
            "title": "Hand-off",
            "detail": ["to BHN / UConn", "clinical team"],
            "watcher": "clinical-lead",
            "modifier": "uncertain-assumption",
            "callout": {"text": "[ASSUMPTION] BHN role unresolved (Cena vs UConn)",
                        "type": "assumption"},
        },
        {
            "id": "audit",
            "x": 580, "y": 420, "w": 200, "h": 80,
            "title": "Audit record sealed",
            "detail": ["chain of custody", "terminal state"],
            "modifier": "default",
        },
    ],
    "edges": [
        {"from": "trigger", "to": "decision", "style": "default"},
        {"from": "decision", "to": "attestation-gate", "label": "red / yellow", "style": "emphasis"},
        {"from": "attestation-gate", "to": "handoff", "label": "signed → handoff", "style": "emphasis"},
        {"from": "handoff", "to": "audit", "label": "ack → audit", "style": "default"},
        {"from": "attestation-gate", "to": "audit", "label": "deferred → audit", "style": "muted"},
        {"from": "handoff", "to": "attestation-gate", "label": "timeout → paged", "style": "muted"},
    ],
    "caption": ("Emphasized path: red/yellow severity → clinical-lead attestation → "
                "BHN/UConn handoff. Audit sealed at every terminal state."),
}


if __name__ == "__main__":
    args = sys.argv[1:]
    if len(args) == 2 and args[0] == "--pl":
        emit_pl(args[1])
    elif len(args) == 3 and args[0] == "--from-md":
        # Parse single use-case markdown → spec → emit one workflow slide
        from spec_from_markdown import spec_from_use_case
        spec = spec_from_use_case(args[1])
        emit_diagram(spec, args[2])
    elif len(args) == 3 and args[0] == "--deck-md":
        # Parse all use-case directories under <root> → emit a multi-slide deck
        from spec_from_markdown import spec_from_use_case
        root = Path(args[1])
        use_case_dirs = sorted(
            d for d in root.iterdir()
            if d.is_dir() and (d / "manifest.md").exists()
        )
        specs = [spec_from_use_case(d) for d in use_case_dirs]
        print(f"  parsed {len(specs)} use-cases from {root}")
        emit_deck(
            specs, args[2],
            deck_title="Cena Workflow Drilldowns",
            deck_subtitle="DRAFT • pending Marrero + Healthcare Data Governance review",
        )
    elif len(args) == 1:
        emit_diagram(ESCALATION_SPEC, args[0])
    else:
        print("Usage:", file=sys.stderr)
        print("  slides-emit.py <out.pptx>                              # hardcoded ESCALATION_SPEC", file=sys.stderr)
        print("  slides-emit.py --pl <out.pptx>                         # Slides Pattern Library", file=sys.stderr)
        print("  slides-emit.py --from-md <use-case-dir> <out.pptx>     # one workflow from markdown", file=sys.stderr)
        print("  slides-emit.py --deck-md <use-cases-root> <out.pptx>   # multi-slide deck from all use-cases", file=sys.stderr)
        sys.exit(64)
